from __future__ import annotations

import json
import os
import re
from datetime import date, datetime, timedelta
from pathlib import Path

from models import Assignment, ExamEvent, StudyGroup, StudyGroupMember, StudyPlan, StudyPlanItem, StudyTopic, Notification, User, Module, db

ALLOWED_EXTENSIONS = {"pdf", "ppt", "pptx", "doc", "docx", "txt", "png", "jpg", "jpeg"}
STATUS_ORDER = ["Not Started", "In Progress", "Submitted", "Completed"]
EMAIL_RE = re.compile(r"^[^@\s]+@[^@\s]+\.[^@\s]+$")
MODULE_CODE_RE = re.compile(r"^[A-Za-z0-9][A-Za-z0-9\- ]{1,19}$")
PASSWORD_MIN_LENGTH = 8
GEMINI_MODEL = os.getenv("GEMINI_MODEL", "gemini-2.5-flash")


def clean_text(value: str | None) -> str:
    return (value or "").strip()


def is_valid_email(email: str) -> bool:
    return bool(EMAIL_RE.match(email or ""))


def is_strong_password(password: str) -> tuple[bool, str]:
    if len(password or "") < PASSWORD_MIN_LENGTH:
        return False, f"Password must be at least {PASSWORD_MIN_LENGTH} characters."
    if not re.search(r"[A-Za-z]", password):
        return False, "Password must include at least one letter."
    if not re.search(r"\d", password):
        return False, "Password must include at least one number."
    return True, ""


def is_valid_module_code(code: str) -> bool:
    return bool(MODULE_CODE_RE.match(code or ""))


def bound_int(value, field_name: str, minimum: int | None = None, maximum: int | None = None):
    try:
        number = int(value)
    except (TypeError, ValueError):
        raise ValueError(f"{field_name} must be a whole number.")
    if minimum is not None and number < minimum:
        raise ValueError(f"{field_name} must be at least {minimum}.")
    if maximum is not None and number > maximum:
        raise ValueError(f"{field_name} must be at most {maximum}.")
    return number


def parse_dt_local(value: str, field_name: str):
    try:
        return datetime.strptime(value, "%Y-%m-%dT%H:%M")
    except (TypeError, ValueError):
        raise ValueError(f"{field_name} must be a valid date and time.")


def parse_date(value: str, field_name: str):
    try:
        return datetime.strptime(value, "%Y-%m-%d").date()
    except (TypeError, ValueError):
        raise ValueError(f"{field_name} must be a valid date.")


def allowed_file(filename: str) -> bool:
    return "." in filename and filename.rsplit(".", 1)[1].lower() in ALLOWED_EXTENSIONS


def cycle_status(current: str) -> str:
    try:
        idx = STATUS_ORDER.index(current)
    except ValueError:
        idx = 0
    return STATUS_ORDER[(idx + 1) % len(STATUS_ORDER)]


def status_class(status: str) -> str:
    mapping = {
        "Not Started": "status-not-started",
        "In Progress": "status-in-progress",
        "Submitted": "status-submitted",
        "Completed": "status-completed",
    }
    return mapping.get(status, "status-not-started")


def calculate_topic_priority(difficulty: int, preference: int) -> int:
    return (difficulty * 2) + (10 - preference)


def _gemini_client():
    api_key = os.getenv("GEMINI_API_KEY") or os.getenv("GOOGLE_API_KEY")
    if not api_key:
        return None
    try:
        from google import genai
        return genai.Client(api_key=api_key)
    except Exception:
        return None


def _parse_difficulty_response(text: str) -> int | None:
    if not text:
        return None
    try:
        data = json.loads(text)
        if isinstance(data, dict):
            for key in ("difficulty", "score"):
                if key in data:
                    return max(1, min(10, int(data[key])))
    except Exception:
        pass
    match = re.search(r"(10|[1-9])", text)
    if match:
        return int(match.group(1))
    return None


def score_topic_with_gemini(title: str, file_text: str = "") -> int | None:
    client = _gemini_client()
    if client is None:
        return None

    prompt = f"""
You are scoring the difficulty of a university study topic for a student planner.
Return JSON only with this exact structure:
{{"difficulty": <integer 1-10>, "reason": "short reason"}}

Topic title: {title}
Document text excerpt: {file_text[:6000]}

Rules:
- 1 means very easy, 10 means very difficult.
- Prefer higher difficulty for dense, conceptual, technical, or long material.
- Keep reason short.
""".strip()

    try:
        response = client.models.generate_content(model=GEMINI_MODEL, contents=prompt)
        text = getattr(response, "text", "") or ""
        difficulty = _parse_difficulty_response(text)
        if difficulty is not None:
            return difficulty
    except Exception:
        return None
    return None


def estimate_difficulty_from_text(title: str, file_text: str = "") -> int:
    gemini_score = score_topic_with_gemini(title, file_text)
    if gemini_score is not None:
        return gemini_score

    score = 4
    if len(title.split()) >= 4:
        score += 1
    title_lower = title.lower()
    if any(k in title_lower for k in ["project", "analysis", "design", "implementation", "proof"]):
        score += 2
    if any(k in title_lower for k in ["quiz", "recap", "intro"]):
        score -= 1
    if file_text:
        words = len(file_text.split())
        if words > 2000:
            score += 2
        elif words > 800:
            score += 1
    return max(1, min(10, score))


def extract_uploaded_text(path: Path) -> str:
    suffix = path.suffix.lower()
    try:
        if suffix == ".txt":
            return path.read_text(encoding="utf-8", errors="ignore")
        if suffix == ".pdf":
            try:
                from pypdf import PdfReader
                reader = PdfReader(str(path))
                return "\n".join((page.extract_text() or "") for page in reader.pages[:20])
            except Exception:
                pass
        if suffix == ".docx":
            try:
                from docx import Document
                doc = Document(str(path))
                return "\n".join(p.text for p in doc.paragraphs)
            except Exception:
                pass
        if suffix == ".pptx":
            try:
                from pptx import Presentation
                prs = Presentation(str(path))
                lines = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            lines.append(shape.text)
                return "\n".join(lines)
            except Exception:
                pass
        return path.read_text(encoding="utf-8", errors="ignore")
    except Exception:
        return ""


def create_exam_anchor(module: Module, exam_start_date: date, exam_end_date: date, assignment_id: int | None = None, exam_kind: str = "mid"):
    existing = ExamEvent.query.filter_by(
        module_id=module.id,
        exam_start_date=exam_start_date,
        exam_end_date=exam_end_date,
        exam_kind=exam_kind,
    ).first()
    if existing:
        return existing
    exam = ExamEvent(
        module_id=module.id,
        module_code=module.code,
        exam_kind=exam_kind,
        exam_start_date=exam_start_date,
        exam_end_date=exam_end_date,
        created_from_assignment_id=assignment_id,
    )
    db.session.add(exam)
    db.session.commit()
    return exam


def _available_days_before_exam(exam_start_date: date) -> list[date]:
    today = date.today()
    if exam_start_date <= today:
        return [today]
    days = []
    cursor = today
    while cursor < exam_start_date:
        days.append(cursor)
        cursor += timedelta(days=1)
    return days or [today]


def generate_study_plan_for_exam(user: User, module_id: int, exam_start_date: date, exam_end_date: date, plan_name: str):
    topics = StudyTopic.query.filter_by(module_id=module_id).order_by(StudyTopic.priority_value.desc(), StudyTopic.id.asc()).all()
    if not topics:
        return None

    plan = StudyPlan(
        user_id=user.id,
        module_id=module_id,
        exam_start_date=exam_start_date,
        exam_end_date=exam_end_date,
        plan_name=plan_name,
        is_shared=False,
    )
    db.session.add(plan)
    db.session.flush()

    study_days = _available_days_before_exam(exam_start_date)
    for index, topic in enumerate(topics):
        target_day = study_days[min(index, len(study_days) - 1)]
        db.session.add(StudyPlanItem(
            plan_id=plan.id,
            topic_id=topic.id,
            target_date=target_day,
            completed=False,
            note=""
        ))

    cursor = exam_start_date
    while cursor <= exam_end_date:
        db.session.add(StudyPlanItem(
            plan_id=plan.id,
            topic_id=None,
            target_date=cursor,
            completed=False,
            note=f"Exam window: {plan_name}"
        ))
        cursor += timedelta(days=1)

    db.session.commit()
    return plan


def reschedule_missed_items(plan: StudyPlan):
    today = date.today()
    items = StudyPlanItem.query.filter_by(plan_id=plan.id).order_by(StudyPlanItem.target_date.asc()).all()
    next_day = today
    for item in items:
        if item.completed or item.topic_id is None:
            continue
        if item.target_date < today:
            item.target_date = next_day
            next_day += timedelta(days=1)
    db.session.commit()


def _group_match_score(user: User, group: StudyGroup) -> int:
    members = StudyGroupMember.query.filter_by(group_id=group.id).all()
    others = [User.query.get(m.user_id) for m in members]
    others = [u for u in others if u and u.id != user.id]
    if not others:
        return 1

    scores = []
    for other in others:
        score = 0
        if user.location_pref and other.location_pref and user.location_pref.lower() == other.location_pref.lower():
            score += 3
        if user.study_goal and other.study_goal and user.study_goal.lower() == other.study_goal.lower():
            score += 2
        if user.availability and other.availability:
            u = set(re.split(r"[\s,;/]+", user.availability.lower().strip()))
            o = set(re.split(r"[\s,;/]+", other.availability.lower().strip()))
            score += len([x for x in u.intersection(o) if x])
        scores.append(score)
    return sum(scores) // len(scores)


def create_or_join_group_for_exam(module: Module, exam_start_date: date, exam_end_date: date, user: User):
    candidates = StudyGroup.query.filter_by(module_code=module.code, exam_start_date=exam_start_date, exam_end_date=exam_end_date, status="open").all()

    chosen = None
    best_score = -1
    for group in candidates:
        score = _group_match_score(user, group)
        if score > best_score:
            best_score = score
            chosen = group

    if chosen is None or best_score < 1:
        suffix = len(candidates) + 1
        chosen = StudyGroup(
            module_id=module.id,
            module_code=module.code,
            exam_start_date=exam_start_date,
            exam_end_date=exam_end_date,
            location_pref=user.location_pref or "",
            group_name=f"{module.code} study circle {suffix}",
            status="open",
        )
        db.session.add(chosen)
        db.session.flush()
    elif not chosen.location_pref and user.location_pref:
        chosen.location_pref = user.location_pref

    existing_member = StudyGroupMember.query.filter_by(group_id=chosen.id, user_id=user.id).first()
    if not existing_member:
        db.session.add(StudyGroupMember(group_id=chosen.id, user_id=user.id, role="member"))
    db.session.commit()
    return chosen


def notify(user_id: int, item_type: str, item_id: int, message: str):
    exists = Notification.query.filter_by(
        user_id=user_id,
        item_type=item_type,
        item_id=item_id,
        message=message
    ).first()
    if not exists:
        db.session.add(Notification(
            user_id=user_id,
            item_type=item_type,
            item_id=item_id,
            message=message
        ))
        db.session.commit()


def assignments_due_for_reminder():
    now = datetime.utcnow()
    items = Assignment.query.filter(Assignment.status != "Completed").all()
    due = []
    for assignment in items:
        delta = assignment.deadline - now
        if timedelta(days=6, hours=23) <= delta <= timedelta(days=7):
            due.append((assignment, "7 days"))
        elif timedelta(hours=23, minutes=59) <= delta <= timedelta(days=1):
            due.append((assignment, "24 hours"))
        elif timedelta(0) <= delta <= timedelta(hours=2):
            due.append((assignment, "2 hours"))
    return due
